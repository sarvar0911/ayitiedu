import os
import tempfile
import subprocess
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime
from io import BytesIO

def generate_certificate(student_name, template_path):
    prs = Presentation(template_path)
    today = datetime.now().strftime('%Y-%m-%d')
    to_date = (datetime.now().replace(year=datetime.now().year + 3)).strftime('%Y-%m-%d')
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if "{{student_name}}" in shape.text:
                    shape.text = shape.text.replace("{{student_name}}", student_name)
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = 'Times New Roman'
                        paragraph.font.size = Pt(32)

                if "{{date}}" in shape.text:
                    shape.text = shape.text.replace("{{date}}", today)
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = 'Times New Roman'
                        paragraph.font.size = Pt(18)
                        paragraph.font.italic = True

                if "{{to_date}}" in shape.text:
                    shape.text = shape.text.replace("{{to_date}}", to_date)
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = 'Times New Roman'
                        paragraph.font.size = Pt(14)
                        paragraph.font.italic = True

    temp_pptx_path = os.path.join(tempfile.gettempdir(), f'temp_certificate_{student_name}.pptx')
    prs.save(temp_pptx_path)

    temp_pdf_path = os.path.join(tempfile.gettempdir(), f'certificate_{student_name}.pdf')
    pdf_stream = convert_pptx_to_pdf(temp_pptx_path, temp_pdf_path)

    cleanup_temp_files([temp_pptx_path, temp_pdf_path])

    return pdf_stream

def convert_pptx_to_pdf(pptx_path, pdf_path):
    try:
        command = ['python', 'C:\\scripts\\unoconv\\unoconv', '-f', 'pdf', '-o', pdf_path, pptx_path]
        subprocess.run(command, check=True)

        with open(pdf_path, 'rb') as pdf_file:
            return BytesIO(pdf_file.read())

    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"Error converting PPTX to PDF: {str(e)}")

def cleanup_temp_files(file_paths):
    for file_path in file_paths:
        if os.path.exists(file_path):
            os.remove(file_path)
